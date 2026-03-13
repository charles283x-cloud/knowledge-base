import os
import uuid
from datetime import datetime

import mammoth
import openpyxl
from pptx import Presentation
from flask import (
    Flask, render_template, request, redirect, url_for,
    flash, send_from_directory, abort, jsonify
)
from flask_login import (
    LoginManager, login_user, logout_user,
    login_required, current_user
)

from config import (
    SECRET_KEY, SQLALCHEMY_DATABASE_URI, UPLOAD_FOLDER,
    MAX_CONTENT_LENGTH, ALLOWED_EXTENSIONS
)
from models import Database, User, Document, Folder, ContactMessage, NewsArticle

app = Flask(__name__)
app.secret_key = SECRET_KEY
app.config['MAX_CONTENT_LENGTH'] = MAX_CONTENT_LENGTH

APP_PREFIX = os.environ.get('APP_PREFIX', '')

class PrefixMiddleware:
    """Strip and set SCRIPT_NAME so Flask generates correct URLs behind a subpath."""
    def __init__(self, wsgi_app, prefix=''):
        self.wsgi_app = wsgi_app
        self.prefix = prefix

    def __call__(self, environ, start_response):
        if self.prefix:
            environ['SCRIPT_NAME'] = self.prefix
            path = environ.get('PATH_INFO', '')
            if path.startswith(self.prefix):
                environ['PATH_INFO'] = path[len(self.prefix):] or '/'
        return self.wsgi_app(environ, start_response)

if APP_PREFIX:
    app.wsgi_app = PrefixMiddleware(app.wsgi_app, prefix=APP_PREFIX)

db_path = SQLALCHEMY_DATABASE_URI.replace('sqlite:///', '')
db = Database(db_path)

os.makedirs(UPLOAD_FOLDER, exist_ok=True)

login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'login'
login_manager.login_message = '请先登录'


@login_manager.user_loader
def load_user(user_id):
    return User.get_by_id(db, int(user_id))


def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def get_file_type(filename):
    ext = filename.rsplit('.', 1)[1].lower()
    type_map = {
        'docx': 'word',
        'xlsx': 'excel', 'xls': 'excel',
        'pdf': 'pdf',
        'pptx': 'ppt', 'ppt': 'ppt',
    }
    return type_map.get(ext, 'unknown')


# ---------- Auth Routes ----------

@app.route('/login', methods=['GET', 'POST'])
def login():
    if current_user.is_authenticated:
        return redirect(url_for('index'))

    if request.method == 'POST':
        username = request.form.get('username', '').strip()
        password = request.form.get('password', '')

        user = User.get_by_username(db, username)
        if user and user.check_password(password):
            login_user(user, remember=True)
            next_page = request.args.get('next')
            return redirect(next_page or url_for('index'))

        flash('用户名或密码错误', 'error')

    return render_template('login.html')


@app.route('/logout')
@login_required
def logout():
    logout_user()
    return redirect(url_for('login'))


# ---------- Main Routes ----------

@app.route('/')
@login_required
def index():
    search = request.args.get('q', '').strip()
    folder_id = request.args.get('folder', None, type=int)

    current_folder = None
    breadcrumbs = []
    subfolders = []

    if search:
        documents = Document.get_all(db, search=search)
    else:
        documents = Document.get_all(db, folder_id=folder_id)
        subfolders = Folder.get_children(db, folder_id)
        if folder_id:
            current_folder = Folder.get_by_id(db, folder_id)
            breadcrumbs = Folder.get_breadcrumbs(db, folder_id)

    all_folders = Folder.get_all(db)

    return render_template('index.html',
        documents=documents, search=search,
        current_folder=current_folder, breadcrumbs=breadcrumbs,
        subfolders=subfolders, all_folders=all_folders)


# ---------- Folder Routes ----------

@app.route('/folder/create', methods=['POST'])
@login_required
def create_folder():
    name = request.form.get('name', '').strip()
    parent_id = request.form.get('parent_id', None, type=int)

    if not name:
        flash('文件夹名称不能为空', 'error')
    else:
        Folder.create(db, name, current_user.id, parent_id)
        flash(f'文件夹 "{name}" 创建成功', 'success')

    if parent_id:
        return redirect(url_for('index', folder=parent_id))
    return redirect(url_for('index'))


@app.route('/folder/rename/<int:folder_id>', methods=['POST'])
@login_required
def rename_folder(folder_id):
    new_name = request.form.get('name', '').strip()
    folder = Folder.get_by_id(db, folder_id)
    if not folder:
        abort(404)
    if not new_name:
        flash('文件夹名称不能为空', 'error')
    else:
        Folder.rename(db, folder_id, new_name)
        flash('文件夹已重命名', 'success')

    if folder.parent_id:
        return redirect(url_for('index', folder=folder.parent_id))
    return redirect(url_for('index'))


@app.route('/folder/delete/<int:folder_id>', methods=['POST'])
@login_required
def delete_folder(folder_id):
    folder = Folder.get_by_id(db, folder_id)
    if not folder:
        abort(404)
    parent_id = folder.parent_id
    Folder.delete(db, folder_id)
    flash('文件夹已删除（文件已移至根目录）', 'success')

    if parent_id:
        return redirect(url_for('index', folder=parent_id))
    return redirect(url_for('index'))


@app.route('/doc/move/<int:doc_id>', methods=['POST'])
@login_required
def move_doc(doc_id):
    doc = Document.get_by_id(db, doc_id)
    if not doc:
        abort(404)
    folder_id = request.form.get('folder_id', None)
    if folder_id == '' or folder_id == 'null':
        folder_id = None
    else:
        folder_id = int(folder_id) if folder_id else None

    Document.move_to_folder(db, doc_id, folder_id)
    flash('文件已移动', 'success')
    return redirect(request.referrer or url_for('index'))


# ---------- Upload ----------

@app.route('/upload', methods=['GET', 'POST'])
@login_required
def upload():
    if request.method == 'POST':
        if 'file' not in request.files:
            flash('没有选择文件', 'error')
            return redirect(request.url)

        file = request.files['file']
        if file.filename == '':
            flash('没有选择文件', 'error')
            return redirect(request.url)

        if not allowed_file(file.filename):
            flash('不支持的文件格式，仅允许 .docx / .xlsx / .xls / .pdf / .pptx', 'error')
            return redirect(request.url)

        original_name = file.filename
        ext = original_name.rsplit('.', 1)[1].lower()
        unique_name = f"{uuid.uuid4().hex}.{ext}"
        file_path = os.path.join(UPLOAD_FOLDER, unique_name)
        file.save(file_path)

        file_size = os.path.getsize(file_path)
        file_type = get_file_type(original_name)
        description = request.form.get('description', '').strip()
        folder_id = request.form.get('folder_id', None, type=int)

        Document.create(db, unique_name, original_name, file_type, file_size,
                        current_user.id, description, folder_id)
        flash('文件上传成功', 'success')

        if folder_id:
            return redirect(url_for('index', folder=folder_id))
        return redirect(url_for('index'))

    folder_id = request.args.get('folder', None, type=int)
    all_folders = Folder.get_all(db)
    return render_template('upload.html', all_folders=all_folders, current_folder_id=folder_id)


# ---------- Document Preview ----------

@app.route('/doc/<int:doc_id>')
@login_required
def preview(doc_id):
    doc = Document.get_by_id(db, doc_id)
    if not doc:
        abort(404)

    file_path = os.path.join(UPLOAD_FOLDER, doc.filename)
    if not os.path.exists(file_path):
        abort(404)

    preview_content = ''
    sheets_data = None
    slides_data = None

    if doc.file_type == 'word':
        try:
            with open(file_path, 'rb') as f:
                result = mammoth.convert_to_html(f)
                preview_content = result.value
        except Exception as e:
            preview_content = f'<p class="text-red-500">文档预览失败: {str(e)}</p>'

    elif doc.file_type == 'excel':
        try:
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            sheets_data = {}
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    rows.append([str(cell) if cell is not None else '' for cell in row])
                sheets_data[sheet_name] = rows
            wb.close()
        except Exception as e:
            preview_content = f'<p class="text-red-500">文档预览失败: {str(e)}</p>'

    elif doc.file_type == 'pdf':
        pass  # PDF rendered via iframe in template

    elif doc.file_type == 'ppt':
        try:
            prs = Presentation(file_path)
            slides_data = []
            for i, slide in enumerate(prs.slides):
                slide_content = {'number': i + 1, 'texts': [], 'title': ''}
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_content['texts'].append(text)
                    if hasattr(shape, 'text') and shape.shape_type and shape == slide.shapes.title:
                        slide_content['title'] = shape.text.strip()
                if not slide_content['title'] and slide_content['texts']:
                    slide_content['title'] = slide_content['texts'][0]
                slides_data.append(slide_content)
        except Exception as e:
            preview_content = f'<p class="text-red-500">PPT 预览失败: {str(e)}</p>'

    return render_template('preview.html', doc=doc,
        preview_content=preview_content, sheets_data=sheets_data,
        slides_data=slides_data)


@app.route('/file/<int:doc_id>')
@login_required
def serve_file(doc_id):
    """Serve file inline (for PDF iframe preview)."""
    doc = Document.get_by_id(db, doc_id)
    if not doc:
        abort(404)
    return send_from_directory(UPLOAD_FOLDER, doc.filename, as_attachment=False, download_name=doc.original_name)


@app.route('/download/<int:doc_id>')
@login_required
def download(doc_id):
    doc = Document.get_by_id(db, doc_id)
    if not doc:
        abort(404)
    return send_from_directory(UPLOAD_FOLDER, doc.filename, as_attachment=True, download_name=doc.original_name)


@app.route('/delete/<int:doc_id>', methods=['POST'])
@login_required
def delete_doc(doc_id):
    doc = Document.get_by_id(db, doc_id)
    if not doc:
        abort(404)

    if not current_user.is_admin and current_user.id != doc.uploaded_by:
        flash('没有权限删除此文件', 'error')
        return redirect(url_for('index'))

    folder_id = doc.folder_id
    filename = Document.delete(db, doc_id)
    if filename:
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        if os.path.exists(file_path):
            os.remove(file_path)

    flash('文件已删除', 'success')
    if folder_id:
        return redirect(url_for('index', folder=folder_id))
    return redirect(url_for('index'))


# ---------- Contact API ----------

@app.route('/api/contact', methods=['POST'])
def contact_api():
    data = request.get_json(silent=True)
    if not data:
        return jsonify({'success': False, 'error': 'Invalid request'}), 400

    required = ['company', 'name', 'email', 'message']
    for field in required:
        if not data.get(field, '').strip():
            return jsonify({'success': False, 'error': f'Missing field: {field}'}), 400

    ContactMessage.create(
        db,
        company=data['company'].strip(),
        name=data['name'].strip(),
        position=data.get('position', '').strip(),
        category=data.get('category', '').strip(),
        phone=data['phone'].strip(),
        email=data['email'].strip(),
        message=data['message'].strip()
    )
    return jsonify({'success': True})


# ---------- News API (public) ----------

@app.route('/api/news')
def news_api():
    articles = NewsArticle.get_published(db)
    result = []
    for a in articles:
        result.append({
            'id': a.id,
            'title_ja': a.title_ja,
            'title_zh': a.title_zh,
            'content_ja': a.content_ja,
            'content_zh': a.content_zh,
            'category': a.category,
            'created_at': a.created_at,
        })
    return jsonify(result)


@app.route('/api/news/<int:article_id>')
def news_detail_api(article_id):
    a = NewsArticle.get_by_id(db, article_id)
    if not a or not a.published:
        return jsonify({'error': 'not found'}), 404
    return jsonify({
        'id': a.id,
        'title_ja': a.title_ja,
        'title_zh': a.title_zh,
        'content_ja': a.content_ja,
        'content_zh': a.content_zh,
        'category': a.category,
        'created_at': a.created_at,
    })


# ---------- News Admin ----------

@app.route('/admin/news')
@login_required
def admin_news():
    if not current_user.is_admin:
        abort(403)
    articles = NewsArticle.get_all(db)
    return render_template('admin_news.html', articles=articles)


@app.route('/admin/news/create', methods=['GET', 'POST'])
@login_required
def create_news():
    if not current_user.is_admin:
        abort(403)
    if request.method == 'POST':
        NewsArticle.create(
            db,
            title_ja=request.form.get('title_ja', '').strip(),
            title_zh=request.form.get('title_zh', '').strip(),
            content_ja=request.form.get('content_ja', '').strip(),
            content_zh=request.form.get('content_zh', '').strip(),
            category=request.form.get('category', 'news'),
            created_by=current_user.id,
            published=request.form.get('published') == '1'
        )
        flash('新闻发布成功', 'success')
        return redirect(url_for('admin_news'))
    return render_template('admin_news_form.html', article=None)


@app.route('/admin/news/edit/<int:article_id>', methods=['GET', 'POST'])
@login_required
def edit_news(article_id):
    if not current_user.is_admin:
        abort(403)
    article = NewsArticle.get_by_id(db, article_id)
    if not article:
        abort(404)
    if request.method == 'POST':
        NewsArticle.update(
            db, article_id,
            title_ja=request.form.get('title_ja', '').strip(),
            title_zh=request.form.get('title_zh', '').strip(),
            content_ja=request.form.get('content_ja', '').strip(),
            content_zh=request.form.get('content_zh', '').strip(),
            category=request.form.get('category', 'news'),
            published=request.form.get('published') == '1'
        )
        flash('新闻已更新', 'success')
        return redirect(url_for('admin_news'))
    return render_template('admin_news_form.html', article=article)


@app.route('/admin/news/delete/<int:article_id>', methods=['POST'])
@login_required
def delete_news(article_id):
    if not current_user.is_admin:
        abort(403)
    NewsArticle.delete(db, article_id)
    flash('新闻已删除', 'success')
    return redirect(url_for('admin_news'))


# ---------- Admin Routes ----------

@app.route('/admin')
@login_required
def admin():
    if not current_user.is_admin:
        flash('需要管理员权限', 'error')
        return redirect(url_for('index'))

    users = User.get_all(db)
    messages = ContactMessage.get_all(db)
    return render_template('admin.html', users=users, messages=messages)


@app.route('/admin/add_user', methods=['POST'])
@login_required
def add_user():
    if not current_user.is_admin:
        abort(403)

    username = request.form.get('username', '').strip()
    password = request.form.get('password', '').strip()
    is_admin = request.form.get('is_admin') == '1'

    if not username or not password:
        flash('用户名和密码不能为空', 'error')
        return redirect(url_for('admin'))

    if len(password) < 4:
        flash('密码长度至少4位', 'error')
        return redirect(url_for('admin'))

    if User.create(db, username, password, is_admin):
        flash(f'用户 {username} 创建成功', 'success')
    else:
        flash(f'用户名 {username} 已存在', 'error')

    return redirect(url_for('admin'))


@app.route('/admin/delete_user/<int:user_id>', methods=['POST'])
@login_required
def delete_user(user_id):
    if not current_user.is_admin:
        abort(403)

    if user_id == current_user.id:
        flash('不能删除自己的账号', 'error')
        return redirect(url_for('admin'))

    User.delete(db, user_id)
    flash('用户已删除', 'success')
    return redirect(url_for('admin'))


@app.route('/admin/message/read/<int:msg_id>', methods=['POST'])
@login_required
def mark_message_read(msg_id):
    if not current_user.is_admin:
        abort(403)
    ContactMessage.mark_read(db, msg_id)
    return redirect(url_for('admin'))


@app.route('/admin/message/delete/<int:msg_id>', methods=['POST'])
@login_required
def delete_message(msg_id):
    if not current_user.is_admin:
        abort(403)
    ContactMessage.delete(db, msg_id)
    flash('留言已删除', 'success')
    return redirect(url_for('admin'))


# ---------- Startup ----------

with app.app_context():
    db.init_db()
    from config import ADMIN_USERNAME, ADMIN_PASSWORD
    if not User.get_by_username(db, ADMIN_USERNAME):
        User.create(db, ADMIN_USERNAME, ADMIN_PASSWORD, is_admin=True)


if __name__ == '__main__':
    app.run(debug=True, port=5000)
